from flask import Flask, request, jsonify
from flasgger import Swagger
from email.mime.image import MIMEImage
import os
import shutil
import time
import win32com.client
import pythoncom
from pptx import Presentation
from googletrans import Translator
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

app = Flask(__name__)
swagger = Swagger(app) 
def traduzir_texto(texto):
    translator = Translator()
    traducao = translator.translate(texto, src='pt', dest='en')
    return traducao.text

def enviar_jpg_por_email(email, caminho_jpg):
    servidor_smtp = "smtp.gmail.com"
    porta_smtp = 587
    email_remetente = os.getenv("SCHWARZ_EMAIL_COMUNICACAO")
    senha_remetente = os.getenv("SCHWARZ_SENHA_EMAIL_COMUNICACAO")

    mensagem = MIMEMultipart()
    mensagem["From"] = email_remetente
    mensagem["To"] = email
    mensagem["Subject"] = "Assinatura de e-mail"

    recado = "Email automático, por favor não responder.\n\n\nOlá! Segue sua assinatura nova com o selo GPTW atualizado!\n\nDúvidas RH fica à disposição."
    mensagem_texto = MIMEText(recado)
    mensagem.attach(mensagem_texto)

    with open(caminho_jpg, "rb") as arquivo_jpg:
        anexo_jpg = MIMEImage(arquivo_jpg.read(), _subtype="jpg")
        anexo_jpg.add_header("Content-Disposition", "attachment", filename="Assinatura.jpg")
        mensagem.attach(anexo_jpg)

    with smtplib.SMTP(servidor_smtp, porta_smtp) as servidor:
        servidor.starttls()
        servidor.login(email_remetente, senha_remetente)
        servidor.send_message(mensagem)

def transformar_em_jpg(caminho_arquivo):
    pythoncom.CoInitialize()
    ppttoJPG = 17
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    deck = powerpoint.Presentations.Open(caminho_arquivo)
    time.sleep(2)
    
    caminho_sem_extensao = os.path.splitext(caminho_arquivo)[0]
    deck.SaveAs(caminho_sem_extensao, ppttoJPG)        
    deck.Close()
    powerpoint.Quit()   
    os.remove(caminho_arquivo)

def processar_assinaturas(data):

    nome = data['nome'].upper()
    cargo_portugues = data['cargo'].upper()
    cargo_ingles = traduzir_texto(cargo_portugues)
    ramal = data['ramal']
    celular = data['celular']
    email = data['email']

    if ramal is None:
        ramal = '8700'

    caminho_raiz = os.getcwd()
    if(celular != None):
        caminho_modelo = os.path.join(caminho_raiz, "Assinatura e-mail Schwarz com Celular.pptx")
    else:
        caminho_modelo = os.path.join(caminho_raiz, "Assinatura e-mail Schwarz.pptx")
    caminho_novo_pptx = os.path.join(caminho_raiz, nome + ".pptx")

    shutil.copy2(caminho_modelo, caminho_novo_pptx)

    presentation = Presentation(caminho_novo_pptx)

    slide = presentation.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if "NOME" in run.text:
                        run.text = nome
                    if "CARGOPT" in run.text:
                        run.text = cargo_portugues
                    if "CARGOIN" in run.text:
                        run.text = cargo_ingles.upper()
                    if "RAMAL" in run.text:
                        run.text = run.text.replace("RAMAL", ramal)
                    if "CELULAR" in run.text:
                        run.text = run.text.replace("CELULAR", celular)

    presentation.save(caminho_novo_pptx)

    transformar_em_jpg(caminho_novo_pptx)

    caminho_arquivo_jpg = caminho_novo_pptx.replace(".pptx", "") + "\Slide1.JPG"

    enviar_jpg_por_email(email, caminho_arquivo_jpg)
    shutil.rmtree(caminho_novo_pptx.replace(".pptx", ""))

@app.route('/gerar_assinatura', methods=['GET'])
def gerar_assinatura():
    """
    Gera uma assinatura de e-mail com base nos dados fornecidos.
    ---
    tags:
      - Assinatura
    parameters:
      - name: body
        in: body
        required: true
        schema:
          type: object
          properties:
            nome:
              type: string
              example: Guilherme Gordiano
            cargo:
              type: string
              example: Estagiário
            ramal:
              type: string
              example: 8781
            celular:
              type: string
              example: +55 11 91234-5678
            email:
              type: string
              example: guilherme.gordiano@schwarz.com.br
    responses:
      200:
        description: Assinatura enviada com sucesso
      500:
        description: Erro no processamento
    """
    try:
        data_json = request.json
        data = {
            "nome": data_json.get('nome'),
            "cargo": data_json.get('cargo'),
            "ramal": data_json.get('ramal'),
            "celular": data_json.get('celular'),
            "email": data_json.get('email')
        }
        processar_assinaturas(data)
        return jsonify({"message": "Assinatura enviada com sucesso!"}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0')