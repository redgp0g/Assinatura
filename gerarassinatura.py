from email.mime.image import MIMEImage
import os
import shutil
import time
import win32com.client
import pyodbc
from pptx import Presentation
from googletrans import Translator
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from dotenv import load_dotenv

load_dotenv()

conn_str = os.getenv("STRING_CONNECTION")

def traduzir_texto(texto, origem, destino):
    translator = Translator()
    traducao = translator.translate(texto, src=origem, dest=destino)
    return traducao.text

def buscar_informacoes_funcionario(id_funcionario):

    conn = pyodbc.connect(conn_str)
    cursor = conn.cursor()

    query = f"SELECT nome, cargo, ramal, email FROM Funcionario WHERE idfuncionario = {id_funcionario}"
    cursor.execute(query)

    # Recuperar os resultados da consulta
    resultado = cursor.fetchone()
    if resultado:
        nome = resultado.nome
        cargo_portugues = resultado.cargo
        cargo_ingles = traduzir_texto(cargo_portugues, 'pt', 'en')
        ramal = resultado.ramal
        email = resultado.email
        return nome, cargo_portugues, cargo_ingles, ramal, email
    else:
        return None

def atualizar_slide_com_informacoes(slide, nome, cargo_portugues, cargo_ingles, ramal):

    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if "NOME" in run.text:
                        run.text = run.text.replace("NOME",  nome)
                    if "CARGOPT" in run.text:
                        run.text = run.text.replace("CARGOPT", cargo_portugues)
                    if "CARGOIN" in run.text:
                        run.text = run.text.replace("CARGOIN", cargo_ingles.upper())
                    if "RAMAL" in run.text:
                        run.text = run.text.replace("RAMAL", ramal)

def enviar_jpg_por_email(email, caminho_jpg):
    servidor_smtp = "smtp.gmail.com"
    porta_smtp = 587
    email_remetente = os.getenv("SMTP_SENHA_REMETENTE")
    senha_remetente = os.getenv("SMTP_SENHA_SENHA")

    # Criar o objeto MIMEMultipart
    mensagem = MIMEMultipart()
    mensagem["From"] = email_remetente
    mensagem["To"] = email
    mensagem["Subject"] = "Assinatura de e-mail"

    recado = "Email automático, por favor não responder.\nOlá! Segue sua assinatura corrigida com o selo GPTW atualizado!\nDúvidas RH fica à disposição."
    mensagem_texto = MIMEText(recado)
    mensagem.attach(mensagem_texto)

    # Anexar o JPG à mensagem de e-mail
    with open(caminho_jpg, "rb") as arquivo_jpg:
        anexo_jpg = MIMEImage(arquivo_jpg.read(), _subtype="jpg")
        anexo_jpg.add_header("Content-Disposition", "attachment", filename="Assinatura.jpg")
        mensagem.attach(anexo_jpg)

    # Enviar o e-mail
    with smtplib.SMTP(servidor_smtp, porta_smtp) as servidor:
        servidor.starttls()
        servidor.login(email_remetente, senha_remetente)
        servidor.send_message(mensagem)


def transformar_em_jpg(caminho_arquivo):
    ppttoJPG = 17
    if caminho_arquivo.endswith(".pptx"):
        try:
            powerpoint = win32com.client.Dispatch("Powerpoint.Application")
            deck = powerpoint.Presentations.Open(caminho_arquivo)
            time.sleep(2)
            deck.SaveAs(caminho_arquivo[:-5], ppttoJPG)
            deck.Close()
            powerpoint.Quit()   
            print('Salvo em JPG')
            os.remove(caminho_arquivo)
        except:
            print('Não foi possível abrir o arquivo')

    elif caminho_arquivo.endswith(".ppt"):
        try:
            powerpoint = win32com.client.Dispatch("Powerpoint.Application")
            deck = powerpoint.Presentations.Open(caminho_arquivo)
            deck.SaveAs(caminho_arquivo[:-4], ppttoJPG)
            deck.Close()
            powerpoint.Quit()
            print('Salvo em JPG')
            os.remove(caminho_arquivo)
        except:
            print('Não foi possível abrir o arquivo')


def excluir_arquivo(caminho_arquivo):
    if os.path.exists(caminho_arquivo):
        os.remove(caminho_arquivo)
        print("Arquivo excluído com sucesso.")
    else:
        print("O arquivo não existe.")

def duplicar_arquivo_pptx(caminho_arquivo,diretorio_destino, nome):
    nome_arquivo = os.path.basename(caminho_arquivo)
    path_novo_arquivo = os.path.join(diretorio_destino, nome + ' ' + nome_arquivo)
    shutil.copy2(caminho_arquivo, path_novo_arquivo)
    return path_novo_arquivo


def processar_assinaturas(id):
    ids_funcionarios = buscar_ids_funcionario()

    caminho_arquivo_original = "C:\Temp\Assinatura e-mail Schwarz.pptx"
    diretorio_destino = "C:\Temp\Assinaturas PDF"

    for id_funcionario in ids_funcionarios:
        informacoes_funcionario = buscar_informacoes_funcionario(id_funcionario)

        if informacoes_funcionario:
            nome, cargo_portugues, cargo_ingles, ramal, email = informacoes_funcionario

            if ramal is None:
                ramal = '8700'

            # Duplicar o arquivo PPT original
            caminho_arquivo_duplicado = duplicar_arquivo_pptx(caminho_arquivo_original, diretorio_destino, nome)

            # Carregar o arquivo PPT duplicado
            presentation = Presentation(caminho_arquivo_duplicado)

            # Atualizar o slide com as informações do funcionário
            novo_slide = presentation.slides[0]
            atualizar_slide_com_informacoes(novo_slide, nome, cargo_portugues, cargo_ingles, ramal)

            # Salvar o PPT atualizado
            presentation.save(caminho_arquivo_duplicado)

            # Converter o PPT em PDF
            transformar_em_jpg(caminho_arquivo_duplicado)

            # Caminho do arquivo JPG
            caminho_arquivo_jpg = caminho_arquivo_duplicado.replace(".pptx", "") + "\Slide1.JPG"

            # Enviar o JPG por e-mail
            enviar_jpg_por_email(email, caminho_arquivo_jpg)
            print("Assinatura enviada com sucesso!")
        else:
            print(f"Informações do funcionário com ID {id_funcionario} não encontradas.")


processar_assinaturas()