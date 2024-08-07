from flask import Flask, request, jsonify
from email.mime.image import MIMEImage
import os
import shutil
import time
import win32com.client
import pyodbc
from pptx import Presentation
from googletrans import Translator
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from dotenv import load_dotenv

load_dotenv()

app = Flask(__name__)

conn_str = os.getenv("STRING_CONNECTION")

def traduzir_texto(texto):
    translator = Translator()
    traducao = translator.translate(texto, src='pt', dest='en')
    return traducao.text

def buscar_informacoes_funcionario(id_funcionario):

    conn = pyodbc.connect(conn_str)
    cursor = conn.cursor()

    query = f"SELECT nome, cargo, ramal, email FROM Funcionario WHERE idfuncionario = {id_funcionario}"
    cursor.execute(query)
    resultado = cursor.fetchone()
    nomeSeparado = resultado.nome.split()
    
    if resultado:
        nome = nomeSeparado[0] + " " + nomeSeparado[-1]
        cargo_portugues = resultado.cargo
        cargo_ingles = traduzir_texto(cargo_portugues)
        ramal = resultado.ramal
        email = resultado.email
        return nome, cargo_portugues, cargo_ingles, ramal, email
    else:
        return None

def atualizar_slide(slide, nome, cargo_portugues, cargo_ingles, ramal):

    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if "NOME" in run.text:
                        run.text = nome
                    if "CARGOPT" in run.text:
                        run.text = cargo_portugues
                    if "CARGOIN" in run.text:
                        run.text = cargo_ingles.upper()
                    if "RAMAL" in run.text:
                        run.text = run.text.replace("RAMAL", ramal)

def enviar_jpg_por_email(email, caminho_jpg):
    servidor_smtp = "smtp.gmail.com"
    porta_smtp = 587
    email_remetente = os.getenv("SMTP_EMAIL_REMETENTE")
    senha_remetente = os.getenv("SMTP_SENHA_REMETENTE")

    mensagem = MIMEMultipart()
    mensagem["From"] = email_remetente
    mensagem["To"] = email
    mensagem["Subject"] = "Assinatura de e-mail"

    recado = "Email automático, por favor não responder.\n\n\nOlá! Segue sua assinatura nova com o selo GPTW atualizado!\n\nDúvidas RH fica à disposição."
    mensagem_texto = MIMEText(recado)
    mensagem.attach(mensagem_texto)

    with open(caminho_jpg, "rb") as arquivo_jpg:
        anexo_jpg = MIMEImage(arquivo_jpg.read(), _subtype="jpg")
        anexo_jpg.add_header("Content-Disposition", "attachment", filename="Assinatura.jpg")
        mensagem.attach(anexo_jpg)

    with smtplib.SMTP(servidor_smtp, porta_smtp) as servidor:
        servidor.starttls()
        servidor.login(email_remetente, senha_remetente)
        servidor.send_message(mensagem)

def transformar_em_jpg(caminho_arquivo):
    ppttoJPG = 17
    if caminho_arquivo.endswith(".pptx"):
        try:
            powerpoint = win32com.client.Dispatch("Powerpoint.Application")
            deck = powerpoint.Presentations.Open(caminho_arquivo)
            time.sleep(2)
            deck.SaveAs(caminho_arquivo[:-5], ppttoJPG)
            deck.Close()
            powerpoint.Quit()   
            print('Salvo em JPG')
            os.remove(caminho_arquivo)
        except Exception as e:
            print(f'Não foi possível abrir o arquivo: {e}')

    elif caminho_arquivo.endswith(".ppt"):
        try:
            powerpoint = win32com.client.Dispatch("Powerpoint.Application")
            deck = powerpoint.Presentations.Open(caminho_arquivo)
            deck.SaveAs(caminho_arquivo[:-4], ppttoJPG)
            deck.Close()
            powerpoint.Quit()
            print('Salvo em JPG')
            os.remove(caminho_arquivo)
        except Exception as e:
            print(f'Não foi possível abrir o arquivo: {e}')

def processar_assinaturas(id):
    informacoes_funcionario = buscar_informacoes_funcionario(id)

    if informacoes_funcionario:
        nome, cargo_portugues, cargo_ingles, ramal, email = informacoes_funcionario

        if ramal is None:
            ramal = '8700'

        caminho_raiz = os.getcwd()
        caminho_modelo = os.path.join(caminho_raiz, "Assinatura e-mail Schwarz.pptx")
        caminho_novo_pptx = os.path.join(caminho_raiz, nome + ".pptx")

        shutil.copy2(caminho_modelo, caminho_novo_pptx)

        presentation = Presentation(caminho_novo_pptx)

        novo_slide = presentation.slides[0]
        atualizar_slide(novo_slide, nome, cargo_portugues, cargo_ingles, ramal)

        presentation.save(caminho_novo_pptx)

        transformar_em_jpg(caminho_novo_pptx)

        if email:
            caminho_arquivo_jpg = caminho_novo_pptx.replace(".pptx", "") + "\Slide1.JPG"

            enviar_jpg_por_email(email, caminho_arquivo_jpg)
            shutil.rmtree(caminho_novo_pptx.replace(".pptx", ""))

@app.route('/gerar_assinatura', methods=['POST'])
def api_processar_assinaturas():
    data = request.json
    id_funcionario = data.get('id')
    if not id_funcionario:
        return jsonify({"error": "ID do funcionário é necessário"}), 400

    try:
        processar_assinaturas(id_funcionario)
        return jsonify({"message": "Processamento de assinatura iniciado com sucesso"}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True,host='0.0.0.0')